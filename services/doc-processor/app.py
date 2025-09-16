#!/usr/bin/env python3
"""
Document processing service - Convert various formats to Markdown
Supports PDF, Word, Excel, PPT, images, with hierarchical splitting
"""

import os
import logging
import tempfile
import subprocess
from pathlib import Path
from typing import Optional, Dict, Any, List

from flask import Flask, request, jsonify
from flask_cors import CORS
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pytesseract
from PIL import Image
import magic
import chardet

# Import new splitters
from splitters import (
    ParentChildSplitter,
    SemanticSplitter,
    RecursiveCharacterSplitter
)

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# 创建Flask应用
app = Flask(__name__)
CORS(app)

# 配置
MAX_FILE_SIZE = int(os.getenv('MAX_FILE_SIZE', '50')) * 1024 * 1024  # MB to bytes
UPLOAD_FOLDER = '/app/uploads'
PROCESSED_FOLDER = '/app/processed'
TEMP_FOLDER = '/app/temp'
OCR_ENABLED = os.getenv('OCR_ENABLED', 'true').lower() == 'true'
OCR_LANGUAGE = os.getenv('OCR_LANGUAGE', 'chi_sim+eng')
PANDOC_ENABLED = os.getenv('PANDOC_ENABLED', 'true').lower() == 'true'

# 支持的文件格式
SUPPORTED_FORMATS = {
    'document': ['.pdf', '.doc', '.docx', '.odt', '.rtf', '.txt', '.md'],
    'spreadsheet': ['.xls', '.xlsx', '.csv', '.ods'],
    'presentation': ['.ppt', '.pptx', '.odp'],
    'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.svg']
}

class DocumentConverter:
    """文档转换器类"""
    
    def __init__(self):
        self.mime = magic.Magic(mime=True)
        
    def convert_to_markdown(self, file_path: str, file_type: str = None) -> str:
        """
        将文件转换为Markdown格式
        
        Args:
            file_path: 文件路径
            file_type: 文件类型（可选）
            
        Returns:
            Markdown格式的文本内容
        """
        if not file_type:
            file_type = self._detect_file_type(file_path)
        
        ext = Path(file_path).suffix.lower()
        
        # 根据文件类型选择转换方法
        if ext == '.md':
            return self._read_markdown(file_path)
        elif ext == '.txt':
            return self._convert_text(file_path)
        elif ext == '.pdf':
            return self._convert_pdf(file_path)
        elif ext in ['.doc', '.docx']:
            return self._convert_word(file_path)
        elif ext in ['.xls', '.xlsx']:
            return self._convert_excel(file_path)
        elif ext in ['.ppt', '.pptx']:
            return self._convert_powerpoint(file_path)
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
            return self._convert_image(file_path)
        elif ext == '.csv':
            return self._convert_csv(file_path)
        else:
            # 尝试使用pandoc进行通用转换
            if PANDOC_ENABLED:
                return self._convert_with_pandoc(file_path)
            else:
                raise ValueError(f"Unsupported file format: {ext}")
    
    def _detect_file_type(self, file_path: str) -> str:
        """检测文件类型"""
        return self.mime.from_file(file_path)
    
    def _read_markdown(self, file_path: str) -> str:
        """读取Markdown文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _convert_text(self, file_path: str) -> str:
        """转换文本文件"""
        # 检测编码
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] or 'utf-8'
        
        # 读取文本
        with open(file_path, 'r', encoding=encoding) as f:
            content = f.read()
        
        # 转换为Markdown（添加代码块以保持格式）
        return f"```text\n{content}\n```"
    
    def _convert_pdf(self, file_path: str) -> str:
        """转换PDF文件"""
        markdown_content = []
        
        try:
            # 首先尝试使用PyPDF2提取文本
            with open(file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        markdown_content.append(f"## Page {page_num}\n")
                        markdown_content.append(text)
                        markdown_content.append("\n---\n")
        except Exception as e:
            logger.warning(f"PyPDF2 extraction failed: {e}")
            
            # 如果PyPDF2失败，尝试使用pdftotext（通过subprocess）
            try:
                result = subprocess.run(
                    ['pdftotext', '-layout', file_path, '-'],
                    capture_output=True,
                    text=True,
                    check=True
                )
                markdown_content = [result.stdout]
            except subprocess.CalledProcessError as e:
                logger.error(f"pdftotext failed: {e}")
                
                # 最后尝试OCR
                if OCR_ENABLED:
                    return self._convert_pdf_with_ocr(file_path)
                else:
                    raise ValueError("PDF text extraction failed")
        
        return '\n'.join(markdown_content)
    
    def _convert_pdf_with_ocr(self, file_path: str) -> str:
        """使用OCR转换PDF"""
        # 将PDF转换为图像，然后OCR
        try:
            # 使用pdftoppm将PDF转换为图像
            temp_dir = tempfile.mkdtemp(dir=TEMP_FOLDER)
            subprocess.run(
                ['pdftoppm', '-png', file_path, f'{temp_dir}/page'],
                check=True
            )
            
            markdown_content = []
            for img_file in sorted(Path(temp_dir).glob('*.png')):
                text = pytesseract.image_to_string(
                    Image.open(img_file),
                    lang=OCR_LANGUAGE
                )
                markdown_content.append(text)
            
            return '\n\n---\n\n'.join(markdown_content)
        except Exception as e:
            logger.error(f"PDF OCR failed: {e}")
            raise
    
    def _convert_word(self, file_path: str) -> str:
        """转换Word文档"""
        doc = Document(file_path)
        markdown_lines = []
        
        for paragraph in doc.paragraphs:
            if paragraph.style.name.startswith('Heading'):
                level = int(paragraph.style.name[-1]) if paragraph.style.name[-1].isdigit() else 1
                markdown_lines.append(f"{'#' * level} {paragraph.text}")
            elif paragraph.text.strip():
                markdown_lines.append(paragraph.text)
            markdown_lines.append('')
        
        # 处理表格
        for table in doc.tables:
            markdown_lines.append(self._table_to_markdown(table))
            markdown_lines.append('')
        
        return '\n'.join(markdown_lines)
    
    def _convert_excel(self, file_path: str) -> str:
        """转换Excel文件"""
        wb = load_workbook(file_path, read_only=True)
        markdown_content = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            markdown_content.append(f"# Sheet: {sheet_name}\n")
            
            # 转换为表格
            rows = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows.append([str(cell) if cell is not None else '' for cell in row])
            
            if rows:
                # 创建Markdown表格
                markdown_content.append(self._create_markdown_table(rows))
            
            markdown_content.append("\n---\n")
        
        return '\n'.join(markdown_content)
    
    def _convert_powerpoint(self, file_path: str) -> str:
        """转换PowerPoint文件"""
        prs = Presentation(file_path)
        markdown_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_content.append(f"# Slide {slide_num}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    markdown_content.append(shape.text)
                    markdown_content.append('')
            
            markdown_content.append("\n---\n")
        
        return '\n'.join(markdown_content)
    
    def _convert_image(self, file_path: str) -> str:
        """转换图像文件（使用OCR）"""
        if not OCR_ENABLED:
            return f"![Image]({Path(file_path).name})\n\n*[Image file - OCR disabled]*"
        
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang=OCR_LANGUAGE)
            
            if text.strip():
                return f"# OCR Extracted Text\n\n{text}"
            else:
                return f"![Image]({Path(file_path).name})\n\n*[No text detected in image]*"
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            return f"![Image]({Path(file_path).name})\n\n*[OCR failed: {str(e)}]*"
    
    def _convert_csv(self, file_path: str) -> str:
        """转换CSV文件"""
        import csv
        
        with open(file_path, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            rows = list(reader)
        
        if not rows:
            return "*[Empty CSV file]*"
        
        return self._create_markdown_table(rows)
    
    def _convert_with_pandoc(self, file_path: str) -> str:
        """使用pandoc进行通用转换"""
        try:
            import pypandoc
            output = pypandoc.convert_file(file_path, 'md')
            return output
        except ImportError:
            logger.error("pypandoc not available, pandoc functionality disabled")
            raise ValueError("Pandoc functionality is disabled")
        except Exception as e:
            logger.error(f"Pandoc conversion failed: {e}")
            raise ValueError(f"Cannot convert file: {e}")
    
    def _table_to_markdown(self, table) -> str:
        """将表格转换为Markdown格式"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        
        return self._create_markdown_table(rows)
    
    def _create_markdown_table(self, rows: list) -> str:
        """创建Markdown表格"""
        if not rows:
            return ""
        
        # 确定列数
        num_cols = max(len(row) for row in rows)
        
        # 标准化行（填充空单元格）
        normalized_rows = []
        for row in rows:
            normalized_row = row + [''] * (num_cols - len(row))
            normalized_rows.append(normalized_row)
        
        # 创建表格
        markdown_lines = []
        
        # 表头
        if normalized_rows:
            markdown_lines.append('| ' + ' | '.join(normalized_rows[0]) + ' |')
            markdown_lines.append('|' + '---|' * num_cols)
            
            # 数据行
            for row in normalized_rows[1:]:
                markdown_lines.append('| ' + ' | '.join(row) + ' |')
        
        return '\n'.join(markdown_lines)

# 创建转换器实例
converter = DocumentConverter()

@app.route('/health', methods=['GET'])
def health_check():
    """健康检查端点"""
    return jsonify({
        'status': 'healthy',
        'service': 'doc-processor',
        'ocr_enabled': OCR_ENABLED,
        'pandoc_enabled': PANDOC_ENABLED
    })

@app.route('/convert', methods=['POST'])
def convert_document():
    """转换文档端点"""
    try:
        # 检查文件
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # 检查文件大小
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        
        if file_size > MAX_FILE_SIZE:
            return jsonify({
                'error': f'File too large. Maximum size is {MAX_FILE_SIZE // (1024*1024)}MB'
            }), 413
        
        # 保存文件
        temp_file = tempfile.NamedTemporaryFile(
            delete=False,
            dir=UPLOAD_FOLDER,
            suffix=Path(file.filename).suffix
        )
        file.save(temp_file.name)
        
        try:
            # 转换文档
            markdown_content = converter.convert_to_markdown(temp_file.name)
            
            # 保存结果
            output_file = Path(PROCESSED_FOLDER) / f"{Path(file.filename).stem}.md"
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return jsonify({
                'status': 'success',
                'markdown': markdown_content,
                'original_filename': file.filename,
                'output_file': str(output_file)
            })
            
        finally:
            # 清理临时文件
            os.unlink(temp_file.name)
            
    except Exception as e:
        logger.error(f"Conversion error: {e}")
        return jsonify({
            'error': str(e),
            'status': 'failed'
        }), 500

@app.route('/batch_convert', methods=['POST'])
def batch_convert():
    """Batch convert documents"""
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No files provided'}), 400
        
        results = []
        for file in files:
            if file.filename:
                # Save file
                temp_file = tempfile.NamedTemporaryFile(
                    delete=False,
                    dir=UPLOAD_FOLDER,
                    suffix=Path(file.filename).suffix
                )
                file.save(temp_file.name)
                
                try:
                    # Convert document
                    markdown_content = converter.convert_to_markdown(temp_file.name)
                    
                    results.append({
                        'filename': file.filename,
                        'status': 'success',
                        'markdown': markdown_content
                    })
                except Exception as e:
                    results.append({
                        'filename': file.filename,
                        'status': 'failed',
                        'error': str(e)
                    })
                finally:
                    os.unlink(temp_file.name)
        
        return jsonify({
            'status': 'completed',
            'results': results
        })
        
    except Exception as e:
        logger.error(f"Batch conversion error: {e}")
        return jsonify({
            'error': str(e),
            'status': 'failed'
        }), 500

@app.route('/convert_hierarchical', methods=['POST'])
def convert_hierarchical():
    """Convert document with hierarchical parent-child splitting"""
    try:
        # Check file
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Get splitting parameters
        strategy = request.form.get('strategy', 'parent_child')
        parent_chunk_size = int(request.form.get('parent_chunk_size', 2000))
        child_chunk_size = int(request.form.get('child_chunk_size', 500))
        parent_overlap = int(request.form.get('parent_overlap', 200))
        child_overlap = int(request.form.get('child_overlap', 50))
        
        # Save file
        temp_file = tempfile.NamedTemporaryFile(
            delete=False,
            dir=UPLOAD_FOLDER,
            suffix=Path(file.filename).suffix
        )
        file.save(temp_file.name)
        
        try:
            # Convert to markdown first
            markdown_content = converter.convert_to_markdown(temp_file.name)
            
            # Extract metadata
            metadata = {
                'filename': file.filename,
                'file_type': Path(file.filename).suffix,
                'conversion_time': os.path.getmtime(temp_file.name)
            }
            
            # Apply hierarchical splitting based on strategy
            if strategy == 'parent_child':
                splitter = ParentChildSplitter(
                    parent_chunk_size=parent_chunk_size,
                    child_chunk_size=child_chunk_size,
                    parent_overlap=parent_overlap,
                    child_overlap=child_overlap
                )
                chunks = splitter.split_document(markdown_content, metadata)
                chunk_tree = splitter.get_chunk_tree(chunks)
                
            elif strategy == 'semantic':
                splitter = SemanticSplitter(
                    max_chunk_size=parent_chunk_size,
                    min_chunk_size=child_chunk_size,
                    similarity_threshold=0.5
                )
                chunks = splitter.split_document(markdown_content, metadata)
                chunk_tree = None
                
            elif strategy == 'recursive':
                splitter = RecursiveCharacterSplitter(
                    chunk_size=parent_chunk_size,
                    chunk_overlap=parent_overlap
                )
                chunks = splitter.split_document(markdown_content, metadata)
                chunk_tree = None
                
            else:
                return jsonify({'error': f'Unknown strategy: {strategy}'}), 400
            
            # Format chunks for response
            formatted_chunks = []
            for chunk in chunks:
                chunk_dict = {
                    'content': chunk.content,
                    'chunk_id': chunk.chunk_id,
                    'metadata': chunk.metadata
                }
                
                # Add parent-child specific fields
                if hasattr(chunk, 'parent_id'):
                    chunk_dict['parent_id'] = chunk.parent_id
                if hasattr(chunk, 'children_ids'):
                    chunk_dict['children_ids'] = chunk.children_ids
                if hasattr(chunk, 'level'):
                    chunk_dict['level'] = chunk.level
                    
                formatted_chunks.append(chunk_dict)
            
            response = {
                'status': 'success',
                'strategy': strategy,
                'total_chunks': len(chunks),
                'chunks': formatted_chunks,
                'original_filename': file.filename,
                'markdown_length': len(markdown_content)
            }
            
            if chunk_tree:
                response['chunk_tree'] = chunk_tree
                
            return jsonify(response)
            
        finally:
            # Clean up temp file
            os.unlink(temp_file.name)
            
    except Exception as e:
        logger.error(f"Hierarchical conversion error: {e}")
        return jsonify({
            'error': str(e),
            'status': 'failed'
        }), 500

@app.route('/split_text', methods=['POST'])
def split_text():
    """Split raw text with hierarchical strategy"""
    try:
        data = request.get_json()
        if not data or 'text' not in data:
            return jsonify({'error': 'No text provided'}), 400
        
        text = data['text']
        strategy = data.get('strategy', 'parent_child')
        metadata = data.get('metadata', {})
        
        # Get splitting parameters
        parent_chunk_size = data.get('parent_chunk_size', 2000)
        child_chunk_size = data.get('child_chunk_size', 500)
        parent_overlap = data.get('parent_overlap', 200)
        child_overlap = data.get('child_overlap', 50)
        
        # Apply splitting based on strategy
        if strategy == 'parent_child':
            splitter = ParentChildSplitter(
                parent_chunk_size=parent_chunk_size,
                child_chunk_size=child_chunk_size,
                parent_overlap=parent_overlap,
                child_overlap=child_overlap
            )
            chunks = splitter.split_document(text, metadata)
            chunk_tree = splitter.get_chunk_tree(chunks)
            
        elif strategy == 'semantic':
            splitter = SemanticSplitter(
                max_chunk_size=parent_chunk_size,
                min_chunk_size=child_chunk_size
            )
            chunks = splitter.split_document(text, metadata)
            chunk_tree = None
            
        elif strategy == 'recursive':
            splitter = RecursiveCharacterSplitter(
                chunk_size=parent_chunk_size,
                chunk_overlap=parent_overlap
            )
            chunks = splitter.split_document(text, metadata)
            chunk_tree = None
            
        else:
            return jsonify({'error': f'Unknown strategy: {strategy}'}), 400
        
        # Format chunks for response
        formatted_chunks = []
        for chunk in chunks:
            chunk_dict = {
                'content': chunk.content,
                'chunk_id': chunk.chunk_id,
                'metadata': chunk.metadata
            }
            
            # Add strategy-specific fields
            if hasattr(chunk, 'parent_id'):
                chunk_dict['parent_id'] = chunk.parent_id
            if hasattr(chunk, 'children_ids'):
                chunk_dict['children_ids'] = chunk.children_ids
            if hasattr(chunk, 'level'):
                chunk_dict['level'] = chunk.level
                
            formatted_chunks.append(chunk_dict)
        
        response = {
            'status': 'success',
            'strategy': strategy,
            'total_chunks': len(chunks),
            'chunks': formatted_chunks,
            'text_length': len(text)
        }
        
        if chunk_tree:
            response['chunk_tree'] = chunk_tree
            
        return jsonify(response)
        
    except Exception as e:
        logger.error(f"Text splitting error: {e}")
        return jsonify({
            'error': str(e),
            'status': 'failed'
        }), 500

if __name__ == '__main__':
    # 确保目录存在
    for folder in [UPLOAD_FOLDER, PROCESSED_FOLDER, TEMP_FOLDER]:
        Path(folder).mkdir(parents=True, exist_ok=True)
    
    # 启动服务
    app.run(host='0.0.0.0', port=5003, debug=False)